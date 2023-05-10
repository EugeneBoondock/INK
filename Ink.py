import os
import string
import sys
import time
from dotenv import load_dotenv
import openai
import prompt_toolkit
import click
import requests
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches
from docx2pdf import convert
from docx2pdf import convert
from pptx.util import Cm
from pptx.util import Inches as PptxInches, Pt, Emu
from colormath.color_objects import sRGBColor
from colormath.color_conversions import convert_color
from pptx.dml.color import RGBColor

load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

load_dotenv('un.env')
api_key = os.getenv('UNSPLASH_API_KEY')

def px_to_emu(pixels):
    return int(pixels * 914400 / 96)

def generate_response(prompt, previous_response="", filename="", save_to_file=False):
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=f"You are INK (Intelligent Narrative Keeper), the document generation assistant. You can generate any type of document in any format. Respond to the following request: {prompt} End of request: ",
        max_tokens=2048,
        n=1,
        stop=None,
        temperature=0.5,
    )
    response_text = response.choices[0].text.strip()
    response_text = response_text.lstrip(string.punctuation).strip()

    if not response_text.startswith(tuple(string.ascii_uppercase)):
        response_text = previous_response + " " + response_text

    if save_to_file:
        filename = prompt_toolkit.prompt("Enter the filename to save to (without extension): ")
        filename, ext = os.path.splitext(filename)
        if ext == "":
            ext = prompt_toolkit.prompt("Enter file extension (e.g. .txt, .pdf, .docx, .pptx): ")

        if ext.lower() == ".pdf":
            # Create a DOCX file using python-docx
            document = Document()
            document.add_paragraph(response_text)
            document.add_page_break()
            document.add_heading('Output', level=1)
            p = document.add_paragraph(response_text)
            p.style = document.styles['Normal']
            p_format = p.paragraph_format
            p_format.space_before = Inches(0.2)
            p_format.space_after = Inches(0.2)
            document.save(f"{filename}.docx")

            # Convert the DOCX file to a PDF file using docx2pdf
            convert(f"{filename}.docx", f"{filename}{ext}")

        elif ext.lower() == ".docx":
            # Create a DOCX file using python-docx
            document = Document()
            document.add_paragraph(response_text)
            document.add_page_break()
            document.add_heading('Output', level=1)
            p = document.add_paragraph(response_text)
            p.style = document.styles['Normal']
            p_format = p.paragraph_format
            p_format.space_before = Inches(0.2)
            p_format.space_after = Inches(0.2)
            document.save(f"{filename}{ext}")

        elif ext.lower() == ".pptx":
            # Create a PPTX file using python-pptx
            # Extract capitalized words from prompt and use as Unsplash query
            query_words = [word.strip(string.punctuation) for word in prompt.split() if word[0].isupper()]
            if query_words:
                query = " ".join(query_words)
            else:
                query = "random"

            # Get a random image from Unsplash using the query
            url = f"https://api.unsplash.com/photos/random?query={query}&client_id={api_key}"
            response = requests.get(url)
            json_data = response.json()
            image_url = json_data['urls']['regular']

            # Create PPTX slides with each paragraph on a separate slide
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = ""  # Remove 'Click to add title' text box
            title = slide.shapes.title
            title.text = 'Output'
            response_text_slides = response_text.split("\n\n")
            for i, slide_text in enumerate(response_text_slides):
                if i > 0:
                    slide = prs.slides.add_slide(title_slide_layout)
        
                # Add the image from Unsplash as the background for the slide
                picture = requests.get(image_url).content
                with open('background.jpg', 'wb') as f:
                    f.write(picture)

                slide_bg = slide.background
                slide_bg.fill.solid()
                srgb_color = sRGBColor(255, 255, 255)
                rgb_color = RGBColor(int(srgb_color.rgb_r), int(srgb_color.rgb_g), int(srgb_color.rgb_b))
                slide_bg.fill.fore_color.rgb = rgb_color  # set white background
                with open('background.jpg', 'rb') as f:
                    slide_bg.fill.background_picture = f

        
                subtitle = slide.placeholders[1]
                subtitle.text = slide_text

                # Center the subtitle vertically on the slide and move it to the top
                subtitle.vertical_anchor = "top"
                subtitle.left = Cm(1)
                subtitle.top = Cm(1.5)  # Move up by 1.5 cm (20 points)
                subtitle.width = Cm(22)
                subtitle.height = Cm(10)

            # Save the PPTX file
            prs.save(f"{filename}.pptx")
            print(f"{filename}.pptx created successfully!")

        else:
            # For other file types, simply write the response_text to the file
            with open(f"{filename}{ext}", "w") as f:
                f.write(response_text)
        click.echo(click.style(f"INK: {filename}{ext} saved", fg="green"))
    return response_text


def exit_chat():
    sys.stdout.write("INK: Stopping Ink Spillage")
    for i in range(3):
        time.sleep(1)
        sys.stdout.write(".")
        sys.stdout.flush()
    print("\n")


filename = ""
save_to_file = False

while True:
    user_input = prompt_toolkit.prompt("User: ")
    if user_input.lower() in ["exit", "quit", "bye", "goodbye"]:
        exit_chat()
        break
    if "write" in user_input.lower():
        save_to_file = True
    response = generate_response(user_input, filename=filename, save_to_file=save_to_file)
    if save_to_file:
        filename = response
    else:
        click.echo(click.style(f"INK: {response}", fg="green"))
    save_to_file = False